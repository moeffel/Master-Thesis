#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
SlideSmith GPT: Automatisches Erstellen von PowerPoint-Präsentationen aus Text.
Version: 1.2 - Verbesserte Segmentierung und Design-Grundlagen
"""

import argparse
import logging
import os
import re
from datetime import datetime
from typing import List, Dict, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

# --- Konfiguration ---
# Farben
COLOR_PRIMARY = RGBColor(0x00, 0x40, 0x80)  # Dunkleres Blau
COLOR_SECONDARY = RGBColor(0x66, 0x66, 0x66) # Mittelgrau
COLOR_ACCENT = RGBColor(0xFFA500)   # Orange
COLOR_BACKGROUND = RGBColor(0xFF, 0xFF, 0xFF) # Weiß
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)      # Dunkelgrau (Text)
COLOR_TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF) # Weiß (Text auf dunklem Grund)
COLOR_FOOTER_TEXT = RGBColor(0x80, 0x80, 0x80) # Hellgrau für Footer

# Typografie
FONT_FAMILY = "Calibri" # Standard-PPT-Schrift, oft verfügbar
FONT_TITLE_SIZE = Pt(36)
FONT_SUBTITLE_SIZE = Pt(20)
FONT_SLIDE_TITLE_SIZE = Pt(30)
FONT_BULLET_SIZE = Pt(18)
FONT_BULLET_LEVEL2_SIZE = Pt(16)
FONT_FOOTER_SIZE = Pt(9)
COMPANY_NAME_FOOTER = "SlideSmith Inc." # Optional für Footer

# Layout & Design
MARGIN_LEFT_RIGHT = Cm(1.8)
MARGIN_TOP_BOTTOM = Cm(1.27) # Standard 0.5 Zoll
MAX_BULLETS_PER_SLIDE = 6
MAX_WORDS_PER_SLIDE_BODY = 50 # Nur für den Bullet-Point-Bereich
TITLE_MAX_WORDS = 10 # Für automatische Titelerkennung

# Ausgabe
OUTPUT_FILENAME_PATTERN = "slidesmith_output_{timestamp}.pptx"

# Logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger("SlideSmith")

# --- Kernlogik ---

class OutlineItem:
    def __init__(self, title: str, subtitle: Optional[str] = None, bullets: Optional[List[str]] = None, is_title_slide: bool = False):
        self.title = title.strip() if title else ""
        self.subtitle = subtitle.strip() if subtitle else None
        self.bullets = [b.strip() for b in bullets if b.strip()] if bullets else []
        self.is_title_slide = is_title_slide

    def __repr__(self) -> str:
        return f"OutlineItem(title='{self.title}', subtitle='{self.subtitle}', bullets_count={len(self.bullets)}, is_title={self.is_title_slide})"

def load_text(filepath: str) -> str:
    logger.info(f"Lese Eingabedatei: {filepath}")
    try:
        with open(filepath, 'r', encoding='utf-8') as f: content = f.read()
        logger.info(f"Datei erfolgreich gelesen ({len(content)} Zeichen).")
        return content
    except FileNotFoundError: logger.error(f"Datei nicht gefunden: {filepath}"); raise
    except IOError as e: logger.error(f"Fehler beim Lesen: {e}"); raise

def _is_potential_title(text: str, max_words: int = TITLE_MAX_WORDS) -> bool:
    """Prüft, ob ein Text ein potenzieller Titel ist (kurz, nicht mit Kleinbuchstaben beginnend)."""
    if not text: return False
    words = text.split()
    if not words: return False
    # Erster Buchstabe des ersten Wortes sollte groß sein (einfache Heuristik)
    # und keine typischen Satzanfänge, die keine Titel sind
    if not words[0][0].isupper():
        return False
    if words[0].lower() in ["und", "oder", "aber", "die", "der", "das", "ein", "eine"]:
        return False
    # Nicht zu lang und endet nicht mit einem typischen Satzzeichen für Fließtext
    return len(words) <= max_words and text[-1] not in ['.', '?', '!']


def parse_outline(text: str) -> List[OutlineItem]:
    logger.info("Beginne mit verbesserter Analyse und Strukturierung des Textes...")
    outline: List[OutlineItem] = []
    
    # Gesamten Text in Absätze aufteilen (doppelte Zeilenumbrüche)
    paragraphs = [p.strip() for p in re.split(r'\n\s*\n', text.strip()) if p.strip()]
    if not paragraphs:
        logger.warning("Keine validen Absätze im Text gefunden.")
        return []

    # 1. Titelfolie (Heuristik: erster kurzer Absatz als Titel, zweiter als Untertitel)
    if paragraphs:
        potential_main_title = ""
        potential_main_subtitle = None
        consumed_paragraphs = 0

        # Versuche, expliziten Titel #Titel zu finden
        first_line_of_first_paragraph = paragraphs[0].split('\n')[0].strip()
        if first_line_of_first_paragraph.startswith("# "):
            potential_main_title = first_line_of_first_paragraph[2:].strip()
            # Rest des ersten Absatzes könnte Untertitel sein
            remaining_lines = paragraphs[0].split('\n')[1:]
            if remaining_lines:
                potential_main_subtitle = "\n".join(s.strip() for s in remaining_lines).strip()
            consumed_paragraphs = 1
        elif _is_potential_title(paragraphs[0], max_words=12): # Erster Absatz als Titel, wenn er titelig aussieht
            potential_main_title = paragraphs[0]
            consumed_paragraphs = 1
            if len(paragraphs) > 1 and len(paragraphs[1].split()) < 20 and not paragraphs[1].startswith("#"): # Zweiter als Untertitel
                potential_main_subtitle = paragraphs[1]
                consumed_paragraphs = 2
        
        if potential_main_title:
            outline.append(OutlineItem(title=potential_main_title, subtitle=potential_main_subtitle, is_title_slide=True))
            logger.info(f"Titelfolie: '{potential_main_title}'" + (f" / '{potential_main_subtitle}'" if potential_main_subtitle else ""))
            paragraphs = paragraphs[consumed_paragraphs:] # Verbrauchte Absätze entfernen
        else:
            logger.info("Kein klarer Haupttitel für die Präsentation gefunden. Erstelle Standard-Titelfolie, falls Inhalt folgt.")
            # Fallback: Erstelle später eine Titelfolie, wenn anderer Inhalt kommt
            # oder wenn der erste "normale" Folientitel der Haupttitel wird.
            # Fürs Erste: Keine explizite Titelfolie, wenn kein Titel erkannt wurde.

    # 2. Inhaltsfolien
    current_slide_title: Optional[str] = None
    current_bullets: List[str] = []

    for para_idx, paragraph_text in enumerate(paragraphs):
        lines_in_para = [line.strip() for line in paragraph_text.split('\n') if line.strip()]
        if not lines_in_para: continue

        first_line_in_para = lines_in_para[0]

        # A. Expliziter Folientitel (# am Anfang des Absatzes oder der ersten Zeile)
        if first_line_in_para.startswith("# "):
            if current_slide_title or current_bullets: # Vorherige Folie abschließen
                outline.append(OutlineItem(title=current_slide_title or "Überblick", bullets=current_bullets))
                logger.debug(f"Folie abgeschlossen: '{current_slide_title or 'Überblick'}'")
            current_slide_title = first_line_in_para[2:].strip()
            current_bullets = []
            # Restliche Zeilen im Absatz als Bullets für diesen Titel
            for line_idx in range(1, len(lines_in_para)):
                if lines_in_para[line_idx].startswith(("* ", "- ")):
                    current_bullets.append(lines_in_para[line_idx][2:].strip())
                else:
                    current_bullets.append(lines_in_para[line_idx])
            continue # Nächster Absatz

        # B. Impliziter Folientitel (erster Satz eines längeren Absatzes als Titel?)
        # Oder wenn der vorherige Folientitel abgeschlossen wurde, wird dieser Absatz zu einer neuen Folie.
        is_new_slide_candidate = False
        if not current_slide_title: # Brauchen einen Titel für die erste Inhaltsfolie
            is_new_slide_candidate = True
        
        # Heuristik: Wenn ein Absatz nicht explizit als Fortsetzung (Bullet) markiert ist,
        # und lang genug ist oder "titelig" wirkt, könnte er eine neue Folie starten.
        if _is_potential_title(first_line_in_para):
             is_new_slide_candidate = True


        if is_new_slide_candidate and (current_slide_title or current_bullets): # Vorherige Folie abschließen
            if not outline or outline[-1].title != (current_slide_title or "Überblick"): # Doppeltes Hinzufügen vermeiden
                outline.append(OutlineItem(title=current_slide_title or "Überblick", bullets=current_bullets))
                logger.debug(f"Folie abgeschlossen: '{current_slide_title or 'Überblick'}'")
            current_slide_title = None # Reset für neuen Titel
            current_bullets = []


        # Titel-Extraktion und Bullet-Point-Extraktion aus dem aktuellen Absatz
        if not current_slide_title and _is_potential_title(first_line_in_para):
            current_slide_title = first_line_in_para
            logger.debug(f"Neuer impliziter Folientitel: '{current_slide_title}'")
            # Restliche Zeilen/Sätze des Absatzes als Bullets
            if len(lines_in_para) > 1:
                for line_idx in range(1, len(lines_in_para)):
                    line_text = lines_in_para[line_idx]
                    if line_text.startswith(("* ", "- ")):
                        current_bullets.append(line_text[2:].strip())
                    else: # Ganze Zeile als Bullet
                        current_bullets.append(line_text)
            # Wenn der Absatz nur aus dem Titel bestand, bleibt current_bullets leer.
        else: # Kein neuer Titel erkannt, Absatzinhalt zu aktuellen Bullets hinzufügen
            if not current_slide_title and not outline: # Erster Inhalt ohne klaren Titel
                 # Versuche, den allerersten Absatz als Haupttitel zu verwenden, falls noch keine Titelfolie da ist
                 if not any(item.is_title_slide for item in outline):
                     outline.append(OutlineItem(title=first_line_in_para, is_title_slide=True))
                     logger.info(f"Fallback: Ersten Absatz als Haupttitel verwendet: '{first_line_in_para}'")
                     current_slide_title = first_line_in_para # Für die erste Inhaltsfolie
                     current_bullets = [] # Reset bullets
                     if len(lines_in_para) > 1:
                         for line_idx in range(1, len(lines_in_para)):
                             current_bullets.append(lines_in_para[line_idx])
                     continue

            # Füge Zeilen des Absatzes als Bullets hinzu
            for line_text in lines_in_para:
                if line_text.startswith(("* ", "- ")):
                    current_bullets.append(line_text[2:].strip())
                else:
                    # Teile längere Zeilen in Sätze auf
                    sentences = re.split(r'(?<=[.!?])\s+', line_text)
                    current_bullets.extend(s.strip() for s in sentences if s.strip())
    
    # Letzte Folie hinzufügen, falls noch Inhalt vorhanden ist
    if current_slide_title or current_bullets:
        outline.append(OutlineItem(title=current_slide_title or "Abschluss", bullets=current_bullets))
        logger.debug(f"Letzte Folie abgeschlossen: '{current_slide_title or 'Abschluss'}'")
    
    # Fallback: Wenn gar keine Folien (außer Titelfolie) erstellt wurden, aber Absätze da waren.
    # Mache jeden Absatz zu einer Folie mit dem ersten Satz als Titel.
    if not outline and paragraphs: # Nur wenn outline komplett leer ist (keine Titelfolie etc.)
        logger.warning("Keine klare Struktur erkannt. Erstelle Folien pro Absatz (einfache Methode).")
        for para_text in paragraphs:
            sents = [s.strip() for s in re.split(r'(?<=[.!?])\s+', para_text) if s.strip()]
            if sents:
                title = sents[0]
                bullets = sents[1:] if len(sents) > 1 else []
                outline.append(OutlineItem(title=title, bullets=bullets))

    # Wenn die erste Folie keine Titelfolie ist, mache sie zur Titelfolie (Notfall)
    if outline and not outline[0].is_title_slide:
        logger.info("Mache erste Inhaltsfolie zur Titelfolie, da keine explizite gefunden wurde.")
        outline[0].is_title_slide = True
        # Ggf. Bullets zu Untertitel umwandeln
        if outline[0].bullets and not outline[0].subtitle:
            outline[0].subtitle = " ".join(outline[0].bullets)
            outline[0].bullets = []


    logger.info(f"Textanalyse abgeschlossen. {len(outline)} potenzielle Folien identifiziert.")
    return outline


def _apply_master_styles(prs: Presentation):
    logger.info("Anwenden von überarbeiteten Master-Folien-Stilen...")
    if not prs.slide_masters: logger.error("Kein Slide Master gefunden."); return
    slide_master = prs.slide_masters[0]

    # Hintergrund
    fill = slide_master.background.fill; fill.solid(); fill.fore_color.rgb = COLOR_BACKGROUND

    # --- Layout 0: Titelfolie ---
    try:
        title_layout = slide_master.slide_layouts[0] # pptx.enum.slide.SLD_LAYOUT_TITLE
        
        # Titel-Platzhalter
        title_ph = _find_placeholder(title_layout.placeholders, PP_PLACEHOLDER.TITLE) or \
                   _find_placeholder(title_layout.placeholders, PP_PLACEHOLDER.CENTER_TITLE) # Oft idx 0
        if title_ph:
            title_ph.left, title_ph.top = MARGIN_LEFT_RIGHT, MARGIN_TOP_BOTTOM + Cm(2)
            title_ph.width = prs.slide_width - 2 * MARGIN_LEFT_RIGHT
            title_ph.height = Cm(3)
            tf = title_ph.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE; tf.word_wrap = True
            p = tf.add_paragraph(); p.text = "Haupttitel"; p.font.name = FONT_FAMILY; p.font.size = FONT_TITLE_SIZE
            p.font.bold = True; p.font.color.rgb = COLOR_PRIMARY; p.alignment = PP_ALIGN.LEFT

        # Untertitel-Platzhalter
        subtitle_ph = _find_placeholder(title_layout.placeholders, PP_PLACEHOLDER.SUBTITLE) # Oft idx 1
        if not subtitle_ph: # Fallback, falls SUBTITLE nicht da ist (z.B. in "Title Only" Layout)
            subtitle_ph = _find_placeholder(title_layout.placeholders, PP_PLACEHOLDER.BODY)

        if subtitle_ph:
            subtitle_ph.left, subtitle_ph.top = MARGIN_LEFT_RIGHT, title_ph.top + title_ph.height + Cm(0.5) if title_ph else MARGIN_TOP_BOTTOM + Cm(5)
            subtitle_ph.width = prs.slide_width - 2 * MARGIN_LEFT_RIGHT
            subtitle_ph.height = Cm(2)
            tf = subtitle_ph.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP; tf.word_wrap = True
            p = tf.add_paragraph(); p.text = "Untertitel"; p.font.name = FONT_FAMILY; p.font.size = FONT_SUBTITLE_SIZE
            p.font.color.rgb = COLOR_SECONDARY; p.alignment = PP_ALIGN.LEFT
        
        # Dekorativer Balken oben
        bar_top = slide_master.shapes.add_shape(MSO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Cm(0.5))
        bar_top.fill.solid(); bar_top.fill.fore_color.rgb = COLOR_PRIMARY
        bar_top.line.fill.background() # Keine Linie

    except IndexError: logger.error("Titelfolien-Layout (0) nicht gefunden.")
    except Exception as e: logger.error(f"Fehler beim Stylen des Titelfolien-Layouts: {e}", exc_info=True)

    # --- Layout 1: Titel und Inhalt ---
    try:
        content_layout = slide_master.slide_layouts[1] # pptx.enum.slide.SLD_LAYOUT_TITLE_AND_CONTENT
        
        # Titel-Platzhalter
        title_ph = _find_placeholder(content_layout.placeholders, PP_PLACEHOLDER.TITLE) # Oft idx 0
        if title_ph:
            title_ph.left, title_ph.top = MARGIN_LEFT_RIGHT, MARGIN_TOP_BOTTOM
            title_ph.width = prs.slide_width - 2 * MARGIN_LEFT_RIGHT
            title_ph.height = Cm(1.5)
            tf = title_ph.text_frame; tf.clear(); tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE; tf.word_wrap = True
            p = tf.add_paragraph(); p.text = "Folientitel"; p.font.name = FONT_FAMILY; p.font.size = FONT_SLIDE_TITLE_SIZE
            p.font.bold = True; p.font.color.rgb = COLOR_PRIMARY; p.alignment = PP_ALIGN.LEFT

        # Inhalts-Platzhalter
        body_ph = _find_placeholder(content_layout.placeholders, PP_PLACEHOLDER.BODY) or \
                  _find_placeholder(content_layout.placeholders, PP_PLACEHOLDER.CONTENT) # Oft idx 1
        if body_ph:
            body_ph.left, body_ph.top = MARGIN_LEFT_RIGHT, (title_ph.top + title_ph.height + Cm(0.5)) if title_ph else MARGIN_TOP_BOTTOM + Cm(2)
            body_ph.width = prs.slide_width - 2 * MARGIN_LEFT_RIGHT
            body_ph.height = prs.slide_height - body_ph.top - MARGIN_TOP_BOTTOM - Cm(1) # Platz für Footer
            tf = body_ph.text_frame; tf.clear(); tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE; tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Cm(0.25) # Innenabstand im Body
            # Level 1 Bullets
            p = tf.paragraphs[0]; p.text = "Erster Punkt"; p.font.name = FONT_FAMILY; p.font.size = FONT_BULLET_SIZE
            p.font.color.rgb = COLOR_TEXT; p.level = 0; p.line_spacing = 1.15
            # Level 2 Bullets
            p2 = tf.add_paragraph(); p2.text = "Unterpunkt"; p2.font.name = FONT_FAMILY; p2.font.size = FONT_BULLET_LEVEL2_SIZE
            p2.font.color.rgb = COLOR_TEXT; p2.level = 1; p2.line_spacing = 1.15; p2.left_indent = Inches(0.5)


    except IndexError: logger.error("Inhaltsfolien-Layout (1) nicht gefunden.")
    except Exception as e: logger.error(f"Fehler beim Stylen des Inhaltsfolien-Layouts: {e}", exc_info=True)
    
    # --- Fußzeile auf allen Master-Folien ---
    for i, layout in enumerate(slide_master.slide_layouts):
        try:
            # Seitenzahl (oft Platzhalter mit Typ SLIDE_NUMBER oder FOOTER)
            # Versuche, existierenden zu formatieren oder neuen hinzuzufügen
            page_num_ph = _find_placeholder(layout.placeholders, PP_PLACEHOLDER.SLIDE_NUMBER) or \
                          _find_placeholder(layout.placeholders, PP_PLACEHOLDER.FOOTER)

            footer_y = prs.slide_height - Cm(0.9)
            footer_h = Cm(0.6)

            if page_num_ph and page_num_ph.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                # Standard Seitenzahl-Platzhalter formatieren
                page_num_ph.left = prs.slide_width - MARGIN_LEFT_RIGHT - Cm(2)
                page_num_ph.top = footer_y
                page_num_ph.width = Cm(2)
                page_num_ph.height = footer_h
                tf = page_num_ph.text_frame; tf.clear()
                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                p.text = "<#>"; p.font.name = FONT_FAMILY; p.font.size = FONT_FOOTER_SIZE
                p.font.color.rgb = COLOR_FOOTER_TEXT; p.alignment = PP_ALIGN.RIGHT
                logger.debug(f"Seitenzahl-Platzhalter auf Layout {i} formatiert.")
            else: # Manuell hinzufügen, wenn kein spezifischer Platzhalter da ist
                pn_box = layout.shapes.add_textbox(
                    prs.slide_width - MARGIN_LEFT_RIGHT - Cm(2), footer_y, Cm(2), footer_h )
                tf = pn_box.text_frame; tf.clear(); p = tf.add_paragraph()
                p.text = "<#>"; p.font.name = FONT_FAMILY; p.font.size = FONT_FOOTER_SIZE
                p.font.color.rgb = COLOR_FOOTER_TEXT; p.alignment = PP_ALIGN.RIGHT
                logger.debug(f"Manuelle Seitenzahl auf Layout {i} hinzugefügt.")

            # Optional: Firmenname links
            if COMPANY_NAME_FOOTER:
                cn_box = layout.shapes.add_textbox(MARGIN_LEFT_RIGHT, footer_y, Cm(5), footer_h)
                tf = cn_box.text_frame; tf.clear(); p = tf.add_paragraph()
                p.text = COMPANY_NAME_FOOTER; p.font.name = FONT_FAMILY; p.font.size = FONT_FOOTER_SIZE
                p.font.color.rgb = COLOR_FOOTER_TEXT; p.alignment = PP_ALIGN.LEFT
                logger.debug(f"Firmenname-Footer auf Layout {i} hinzugefügt.")

        except Exception as e:
            logger.warning(f"Fehler beim Hinzufügen/Formatieren der Fußzeile auf Layout {i}: {e}")

    logger.info("Master-Folien-Stile angewendet.")


def build_deck(outline: List[OutlineItem]) -> Tuple[Optional[Presentation], List[OutlineItem]]:
    if not outline: logger.warning("Keine Gliederungsdaten."); return None, []
    prs = Presentation()
    try: _apply_master_styles(prs)
    except Exception as e: logger.error(f"Kritischer Fehler in _apply_master_styles: {e}. Nutze Standard-Layouts.", exc_info=True)
    
    title_slide_layout = prs.slide_layouts[0] if prs.slide_layouts else None
    content_slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else title_slide_layout
    if not title_slide_layout or not content_slide_layout:
        logger.critical("Keine Folienlayouts verfügbar. Abbruch."); return None, []

    appendix_items: List[OutlineItem] = []; slide_count = 0
    logger.info("Beginne Folienerstellung...")

    for item_idx, item in enumerate(outline):
        current_slide_layout = title_slide_layout if item.is_title_slide else content_slide_layout
        
        # Aufteilen von Bullets auf mehrere Folien, falls nötig
        bullets_to_process = list(item.bullets)
        item_title = item.title
        item_subtitle = item.subtitle # Nur für die erste Titelfolie relevant
        page_num_for_item = 0

        while True: # Schleife für Folienaufteilung
            page_num_for_item += 1
            slide_title_text = item_title
            if page_num_for_item > 1 and not item.is_title_slide: # Nur bei Fortsetzungs-Inhaltsfolien
                slide_title_text += f" (Teil {page_num_for_item})"

            slide_bullets = []
            words_on_slide = len(slide_title_text.split())
            
            if not item.is_title_slide: # Nur für Inhaltsfolien die Bullet-Logik anwenden
                temp_remaining_bullets = []
                for b_idx, bullet in enumerate(bullets_to_process):
                    bullet_words = len(bullet.split())
                    if len(slide_bullets) < MAX_BULLETS_PER_SLIDE and (words_on_slide + bullet_words) <= MAX_WORDS_PER_SLIDE_BODY:
                        slide_bullets.append(bullet)
                        words_on_slide += bullet_words
                    else:
                        temp_remaining_bullets.append(bullet)
                bullets_to_process = temp_remaining_bullets
            
            # Nur Folie erstellen, wenn es die erste Folie für das Item ist ODER es Inhalt gibt
            if page_num_for_item == 1 or slide_bullets or (item.is_title_slide and page_num_for_item == 1):
                try:
                    slide = prs.slides.add_slide(current_slide_layout); slide_count += 1
                    logger.info(f"Erstelle Folie {slide_count}: '{slide_title_text}' (Typ: {'Titel' if item.is_title_slide else 'Inhalt'})")

                    # Titel
                    title_ph = _find_placeholder(slide.placeholders, PP_PLACEHOLDER.TITLE) or \
                               _find_placeholder(slide.placeholders, PP_PLACEHOLDER.CENTER_TITLE)
                    if title_ph:
                        tf = title_ph.text_frame; tf.clear(); p = tf.add_paragraph(); p.text = slide_title_text
                        # Stile sollten vom Master kommen, ggf. hier Overrides
                        if item.is_title_slide: p.font.size = FONT_TITLE_SIZE; p.alignment = PP_ALIGN.LEFT
                        else: p.font.size = FONT_SLIDE_TITLE_SIZE; p.alignment = PP_ALIGN.LEFT
                    else: logger.warning(f"Kein Titel-Platzhalter auf Folie {slide_count}.")

                    # Untertitel (nur für Titelfolie)
                    if item.is_title_slide and item_subtitle:
                        subtitle_ph = _find_placeholder(slide.placeholders, PP_PLACEHOLDER.SUBTITLE) or \
                                      _find_placeholder(slide.placeholders, PP_PLACEHOLDER.BODY)
                        if subtitle_ph and subtitle_ph != title_ph: # Sicherstellen, dass es nicht derselbe PH ist
                            tf = subtitle_ph.text_frame; tf.clear(); p = tf.add_paragraph(); p.text = item_subtitle
                        elif item_subtitle: logger.warning(f"Kein separater Untertitel-Platzhalter auf Titelfolie für: {item_subtitle}")
                    
                    # Bullets (nur für Inhaltsfolien)
                    if not item.is_title_slide and slide_bullets:
                        body_ph = _find_placeholder(slide.placeholders, PP_PLACEHOLDER.BODY) or \
                                  _find_placeholder(slide.placeholders, PP_PLACEHOLDER.CONTENT)
                        if body_ph:
                            tf = body_ph.text_frame; tf.clear()
                            for bullet_text in slide_bullets:
                                p = tf.add_paragraph(); p.text = bullet_text # Master sollte Bullet-Stil setzen
                                # Keyword-Hinweise (optional)
                                if any(k in bullet_text.lower() for k in ["diagramm", "chart", "graph"]):
                                     run = p.add_run(); run.text = " ✍️"; run.font.size = Pt(14) # Emoji
                                elif any(k in bullet_text.lower() for k in ["foto", "bild", "image", "picture"]):
                                     run = p.add_run(); run.text = " 🖼️"; run.font.size = Pt(14)
                        else: logger.warning(f"Kein Body-Platzhalter auf Folie {slide_count} für Bullets.")
                    
                    # Wenn Titelfolie, aber kein Untertitel-PH und Bullets vorhanden waren (aus parse_outline), Body-PH entfernen
                    if item.is_title_slide and not item_subtitle:
                        body_ph_on_title_slide = _find_placeholder(slide.placeholders, PP_PLACEHOLDER.BODY)
                        if body_ph_on_title_slide and body_ph_on_title_slide != title_ph:
                             try: body_ph_on_title_slide.element.getparent().remove(body_ph_on_title_slide.element)
                             except: pass # Ignoriere, wenn nicht entfernbar

                except Exception as e:
                    logger.error(f"Fehler beim Erstellen von Folie {slide_count} ('{slide_title_text}'): {e}", exc_info=True)
            
            if not bullets_to_process or item.is_title_slide: # Alle Bullets verarbeitet oder es ist eine Titelfolie
                break # Verlasse die Schleife für Folienaufteilung

        # Verbleibende Bullets in Appendix verschieben
        if bullets_to_process:
            logger.warning(f"{len(bullets_to_process)} Bullets für '{item_title}' passen nicht, -> Appendix.")
            appendix_items.append(OutlineItem(title=f"{item_title} (Forts. Appendix)", bullets=bullets_to_process))

    # --- Appendix Folien ---
    if appendix_items:
        logger.info(f"Erstelle {len(appendix_items)} Appendix-Folie(n)...")
        # Appendix Titelfolie
        if title_slide_layout:
            try:
                slide = prs.slides.add_slide(title_slide_layout); slide_count += 1
                title_ph = _find_placeholder(slide.placeholders, PP_PLACEHOLDER.TITLE) or _find_placeholder(slide.placeholders, PP_PLACEHOLDER.CENTER_TITLE)
                if title_ph: title_ph.text_frame.text = "Appendix"
                # Entferne Untertitel oder Body auf Appendix-Titelfolie
                for ph_type_to_remove in [PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY]:
                    ph_to_remove = _find_placeholder(slide.placeholders, ph_type_to_remove)
                    if ph_to_remove and ph_to_remove != title_ph:
                        try: ph_to_remove.element.getparent().remove(ph_to_remove.element)
                        except: pass
            except Exception as e: logger.error(f"Konnte Appendix-Trennfolie nicht erstellen: {e}")
        
        # Appendix Inhaltsfolien
        for app_item in appendix_items:
            try:
                slide = prs.slides.add_slide(content_slide_layout); slide_count += 1
                title_ph = _find_placeholder(slide.placeholders, PP_PLACEHOLDER.TITLE)
                if title_ph: title_ph.text_frame.text = app_item.title
                body_ph = _find_placeholder(slide.placeholders, PP_PLACEHOLDER.BODY) or _find_placeholder(slide.placeholders, PP_PLACEHOLDER.CONTENT)
                if body_ph:
                    tf = body_ph.text_frame; tf.clear()
                    for bullet in app_item.bullets:
                        p = tf.add_paragraph(); p.text = bullet # Master sollte Stil setzen
                logger.info(f"Appendix-Folie {slide_count}: '{app_item.title}' erstellt.")
            except Exception as e: logger.error(f"Konnte Appendix-Folie '{app_item.title}' nicht erstellen: {e}")

    logger.info(f"Präsentationserstellung abgeschlossen ({slide_count} Folien).")
    return prs, appendix_items


def save_pptx(prs: Presentation, output_dir: str = ".") -> Optional[str]:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = OUTPUT_FILENAME_PATTERN.format(timestamp=timestamp)
    filepath = os.path.join(output_dir, filename)
    try:
        os.makedirs(output_dir, exist_ok=True)
        prs.save(filepath)
        logger.info(f"Präsentation gespeichert: {filepath}")
        return filepath
    except Exception as e:
        logger.error(f"Fehler beim Speichern nach {filepath}: {e}", exc_info=True)
        return None

def run_self_test(prs: Presentation, filepath: str):
    if not prs: logger.error("Test: Kein Präsentationsobjekt."); return
    if not filepath or not os.path.exists(filepath): logger.error(f"Test: Datei '{filepath}' fehlt."); return
    logger.info("Führe Selbsttest durch...")
    passed = True
    try:
        assert os.path.exists(filepath), f"Datei {filepath} existiert nicht."
        logger.info(f"  [OK] Datei existiert: {filepath}")
        assert len(prs.slides) > 0, "Keine Folien."
        logger.info(f"  [OK] Folien: {len(prs.slides)} > 0")
        titles_filled = 0
        for slide in prs.slides:
            title_ph = _find_placeholder(slide.placeholders, PP_PLACEHOLDER.TITLE) or \
                       _find_placeholder(slide.placeholders, PP_PLACEHOLDER.CENTER_TITLE)
            if title_ph and title_ph.has_text_frame and title_ph.text_frame.text.strip(): titles_filled += 1
        assert titles_filled > 0, "Kein Folientitel gefüllt."
        logger.info(f"  [OK] Titel gefüllt: ({titles_filled} gefunden).")
    except AssertionError as e: logger.error(f"Selbsttest fehlgeschlagen: {e}"); passed = False
    except Exception as e: logger.error(f"Fehler im Selbsttest: {e}", exc_info=True); passed = False
    if passed: logger.info("Selbsttest erfolgreich.")
    else: logger.warning("Selbsttest mit Fehlern.")

def main():
    parser = argparse.ArgumentParser(description="SlideSmith GPT", formatter_class=argparse.ArgumentDefaultsHelpFormatter)
    parser.add_argument("input_file", help="Eingabe-Textdatei (.txt).")
    parser.add_argument("-o", "--output-dir", default=".", help="Ausgabeverzeichnis.")
    parser.add_argument("-v", "--verbose", action="store_true", help="Debug-Logs anzeigen.")
    args = parser.parse_args()

    log_level = logging.DEBUG if args.verbose else logging.INFO
    # force=True ist wichtig, wenn main() mehrfach aufgerufen wird (z.B. in Tests)
    # oder wenn andere Bibliotheken das Logging bereits konfiguriert haben.
    logging.basicConfig(level=log_level, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s', force=True)
    logger.info(f"Logging-Level: {logging.getLevelName(logger.getEffectiveLevel())}")

    try:
        raw_text = load_text(args.input_file)
        if not raw_text.strip(): logger.warning("Eingabedatei leer."); return
        outline = parse_outline(raw_text)
        if not outline: logger.warning("Keine Struktur extrahiert. Ggf. Text zu kurz oder unstrukturiert?"); return
        presentation, appendix = build_deck(outline)
        if not presentation: logger.error("Präsentation nicht erstellt."); return
        saved_filepath = save_pptx(presentation, args.output_dir)
        if not saved_filepath: return
        run_self_test(presentation, saved_filepath)
    except FileNotFoundError: print(f"FEHLER: Datei '{args.input_file}' nicht gefunden.")
    except Exception as e: logger.exception(f"Unerwarteter Hauptfehler: {e}"); print(f"Unerwarteter Fehler. Siehe Log.")

if __name__ == "__main__":
    main()